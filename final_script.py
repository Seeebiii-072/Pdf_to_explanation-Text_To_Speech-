# import os
# import io
# import re
# import sqlite3
# import requests
# import time
# # import numpy as np
# from pathlib import Path
# from bs4 import BeautifulSoup
# from time import sleep
#
# # file-format libs
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation
# from PIL import Image
# import pandas as pd
#
# # NLP / tools
# import nltk
# nltk.download('punkt', quiet=True)
# # from nltk.tokenize import sent_tokenize
#
# from sentence_transformers import SentenceTransformer
# from deep_translator import GoogleTranslator
# import pyttsx3
# from gtts import gTTS
# import pytesseract
#
# # NEW MODEL (your model)
# from gpt4all import GPT4All
#
# # === CONFIG ===
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
# OUTPUT_DIR = "output"
# FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
# URL_FETCH_TIMEOUT = 10
# URL_SNIPPET_CHARS = 1200
#
# # NEW MODEL CONFIG
# MODEL_URL = "https://huggingface.co/Qwen/Qwen2.5-1.5B-Instruct-GGUF/resolve/main/qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_FILE = "qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_PATH = f"{MODEL_FOLDER}/{MODEL_FILE}"
# CHUNK_SIZE = 600
#
# DEEP_PROMPT = """
# You are an expert who provides deep and accurate explanations.
#
# You MUST follow these rules exactly:
#
# 1. Only explain the meaning of the given text.
# 2. Do NOT invent, add, assume, or hallucinate any details.
# 3. Do NOT repeat sentences from the text.
# 4. Do NOT summarize.
# 5. Do NOT fix the structure.
# 6. Only explain what exists inside the provided text.
# 7. If the text is unclear, explain ONLY what can be understood, do not guess.
# 8. Explanation must be a single coherent paragraph.
#
# TEXT STARTS:
# {chunk}
# TEXT ENDS.
# """
#
# # Create folders
# os.makedirs(OUTPUT_DIR, exist_ok=True)
# os.makedirs(FETCHED_DIR, exist_ok=True)
# os.makedirs(MODEL_FOLDER, exist_ok=True)
#
# # === LOAD EMBEDDING MODEL ===
# print("🔄 Loading embedding model...")
# embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
# print("✅ Embedding model loaded.")
#
# # === ENSURE NEW MODEL ===
# def ensure_qwen_model():
#     if not os.path.exists(MODEL_PATH):
#         print("📥 Downloading Qwen explanation model...")
#         with requests.get(MODEL_URL, stream=True) as r:
#             with open(MODEL_PATH, "wb") as f:
#                 for chunk in r.iter_content(1024 * 1024):
#                     f.write(chunk)
#         print("✔ Download complete.")
#     else:
#         print("✔ Qwen model exists.")
#
# def load_qwen():
#     print("⏳ Loading Qwen model...")
#     llm = GPT4All(model_name=MODEL_FILE, model_path=MODEL_FOLDER)
#     print("✔ Qwen model loaded.")
#     return llm
#
# ensure_qwen_model()
# qwen = load_qwen()
#
# # === UTILITIES ===
# def clean_path_input(s: str) -> str:
#     return s.strip().replace("\u202a", "").replace("\u202b", "").replace("\ufeff", "")
#
# def read_text_from_file(file_path: str) -> str:
#     ext = Path(file_path).suffix.lower()
#     text = ""
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             text = "\n".join([page.get_text("text") for page in doc])
#             doc.close()
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#         elif ext == ".pptx":
#             text = extract_text_from_pptx(file_path)
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#         elif ext in (".html", ".htm"):
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#         else:
#             raise ValueError(f"Unsupported file type: {ext}")
#     except Exception as e:
#         raise RuntimeError(f"Error reading {file_path}: {e}")
#     return text or ""
#
# def extract_text_from_pptx(pptx_path: str) -> str:
#     prs = Presentation(pptx_path)
#     all_text = []
#     for idx, slide in enumerate(prs.slides, start=1):
#         slide_text = []
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 slide_text.append(shape.text.strip())
#             try:
#                 if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
#                     blob = shape.image.blob
#                     img = Image.open(io.BytesIO(blob)).convert("RGB")
#                     ocr_text = pytesseract.image_to_string(img)
#                     if ocr_text.strip():
#                         slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
#             except:
#                 pass
#         if slide_text:
#             all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
#     return "\n\n".join(all_text)
#
# URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)
#
# def fetch_url_snippet(url: str) -> str:
#     headers = {"User-Agent": "Mozilla/5.0"}
#     try:
#         r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
#         r.raise_for_status()
#         soup = BeautifulSoup(r.text, "html.parser")
#         title = (soup.title.string.strip() if soup.title else "")
#         paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
#         snippet = (title + "\n\n" + " ".join(paragraphs))[:URL_SNIPPET_CHARS]
#         safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
#         with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
#             f.write(r.text)
#         return snippet or title or url
#     except:
#         return url
#
# def expand_links_in_text(text: str) -> str:
#     urls = list(dict.fromkeys(URL_RE.findall(text)))
#     if not urls:
#         return text
#     expanded_text = text
#     for url in urls:
#         snippet = fetch_url_snippet(url)
#         replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
#         expanded_text = expanded_text.replace(url, replacement)
#         sleep(0.2)
#     return expanded_text
#
# def clean_text(text: str) -> str:
#     text = re.sub(r'\r\n', '\n', text)
#     text = re.sub(r'\n{3,}', '\n\n', text)
#     text = re.sub(r'[ \t]+', ' ', text)
#     return text.strip()
#
# # === NEW EXPLANATION USING QWEN ===
# def chunk_words(text, size=CHUNK_SIZE):
#     words = text.split()
#     for i in range(0, len(words), size):
#         yield " ".join(words[i:i + size])
#
# def explain_chunk(chunk):
#     prompt = DEEP_PROMPT.format(chunk=chunk)
#     resp = qwen.generate(prompt, max_tokens=500)
#     return resp.strip()
#
# def generate_deep_explanation(text):
#     final = []
#     for chunk in chunk_words(text):
#         final.append(explain_chunk(chunk))
#     return "\n\n".join(final)
#
# # === DATABASE / EMBEDDINGS ===
# def init_db():
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     c.execute("""
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB,
#             source_file TEXT
#         )
#     """)
#     conn.commit()
#     conn.close()
#
# def save_embeddings(source_file, chunks, vectors):
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     for ch, vec in zip(chunks, vectors):
#         c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
#                   (source_file, ch, vec.tobytes()))
#     conn.commit()
#     conn.close()
#
# # === AUDIO ===
# def save_english_audio(text, out_path):
#     try:
#         tts = pyttsx3.init()
#         tts.setProperty('rate', 150)
#         tts.setProperty('volume', 1.0)
#         tts.save_to_file(text, out_path)
#         tts.runAndWait()
#     except Exception as e:
#         print("❌ English audio failed:", e)
#
# def save_urdu_audio(text, out_path):
#     try:
#         gTTS(text=text.replace("\n", " "), lang='ur').save(out_path)
#     except Exception as e:
#         print("❌ Urdu audio failed:", e)
#
# # === SAFE URDU TRANSLATION ===
# def safe_translate_to_urdu(text, max_chunk=2500, retries=5, delay=2):
#     urdu_output = []
#     parts = [text[i:i+max_chunk] for i in range(0, len(text), max_chunk)]
#     for idx, part in enumerate(parts, start=1):
#         attempt = 1
#         while attempt <= retries:
#             try:
#                 print(f"🌐 Translating chunk {idx}/{len(parts)} (Attempt {attempt})")
#                 translated = GoogleTranslator(source='auto', target='ur').translate(part)
#                 urdu_output.append(translated)
#                 break
#             except Exception as e:
#                 print(f"⚠️ Translation error: {e}")
#                 attempt += 1
#                 time.sleep(delay)
#         if attempt > retries:
#             print("❌ Failed to translate this chunk after multiple retries.")
#             urdu_output.append(part)
#     return " ".join(urdu_output)
#
# # === MAIN PROCESS ===
# def process_document(file_path: str, expand_links: bool = True):
#     print(f"\n📄 Processing: {file_path}")
#
#     raw_text = read_text_from_file(file_path)
#
#     if expand_links:
#         raw_text = expand_links_in_text(raw_text)
#
#     cleaned = clean_text(raw_text)
#     base = Path(file_path).stem
#
#     # Save extracted text
#     raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
#     with open(raw_out, "w", encoding="utf-8") as f:
#         f.write(cleaned)
#
#     # Embeddings
#     chunks = [cleaned[i:i+1000] for i in range(0, len(cleaned), 1000)]
#     vectors = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(file_path, chunks, vectors)
#
#     # NEW DEEP EXPLANATION
#     print("🧠 Generating deep explanation using Qwen...")
#     explanation = generate_deep_explanation(cleaned)
#
#     explain_path = os.path.join(OUTPUT_DIR, f"{base}_deep_explanation.txt")
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation)
#
#     # English audio
#     save_english_audio(explanation, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))
#
#     # Urdu translation
#     ur_summary = safe_translate_to_urdu(explanation)
#
#     # Save Urdu text
#     ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
#     with open(ur_path, "w", encoding="utf-8") as f:
#         f.write(ur_summary)
#
#     # Urdu audio
#     save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))
#
#     print("\n✅ All outputs saved.")
#
# # === CLI ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
#     while True:
#         inp = input("👉 Enter file path: ").strip()
#         inp = clean_path_input(inp)
#         if inp.lower() in ("exit", "quit"):
#             break
#         if not os.path.exists(inp):
#             print("❌ File not found.")
#             continue
#         try:
#             process_document(inp, expand_links=True)
#         except Exception as e:
#             print("❌ Error:", e)

# full working but equation samll explanation
# import os
# import io
# import re
# import sqlite3
# import requests
# import time
# from pathlib import Path
# from bs4 import BeautifulSoup
# from time import sleep
#
# # file-format libs
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation
# from PIL import Image
# import pandas as pd
#
# # NLP / tools
# import nltk
# nltk.download('punkt', quiet=True)
# from sentence_transformers import SentenceTransformer
# from deep_translator import GoogleTranslator
# import pyttsx3
# from gtts import gTTS
# import pytesseract
#
# # NEW MODEL (your model)
# from gpt4all import GPT4All
#
# # === CONFIG ===
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
# OUTPUT_DIR = "output"
# FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
# URL_FETCH_TIMEOUT = 10
# URL_SNIPPET_CHARS = 1200
#
# # NEW MODEL CONFIG
# MODEL_URL = "https://huggingface.co/Qwen/Qwen2.5-1.5B-Instruct-GGUF/resolve/main/qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_FILE = "qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_PATH = f"{MODEL_FOLDER}/{MODEL_FILE}"
# CHUNK_SIZE = 600
#
# DEEP_PROMPT = """
# You are an expert who provides deep and accurate explanations.
#
# You MUST follow these rules exactly:
#
# 1. Only explain the meaning of the given text.
# 2. Do NOT invent, add, assume, or hallucinate any details.
# 3. Do NOT repeat sentences from the text.
# 4. Do NOT summarize.
# 5. Do NOT fix the structure.
# 6. Only explain what exists inside the provided text.
# 7. If the text is unclear, explain ONLY what can be understood, do not guess.
# 8. Explanation must be a single coherent paragraph.
#
# TEXT STARTS:
# {chunk}
# TEXT ENDS.
# """
#
# # General equation prompt for any subject
# EQUATION_PROMPT = """
# You are an expert in explaining formulas and equations from any subject: mathematics, physics, chemistry, economics, or engineering.
#
# Follow these rules exactly:
# 1. Explain ONLY what is present in the equation.
# 2. Describe the meaning of each variable/symbol.
# 3. Explain relationships between variables.
# 4. Explain the usage or context of the formula if apparent.
# 5. Do NOT invent extra unrelated information.
# 6. Make the explanation clear and easy to understand.
#
# EQUATION START:
# {equation}
# EQUATION END.
# """
#
# # Regex to detect inline ($...$) or block ($$...$$) equations
# EQUATION_RE = re.compile(r'(\$\$?.+?\$\$?)')
#
# # Create folders
# os.makedirs(OUTPUT_DIR, exist_ok=True)
# os.makedirs(FETCHED_DIR, exist_ok=True)
# os.makedirs(MODEL_FOLDER, exist_ok=True)
#
# # === LOAD EMBEDDING MODEL ===
# print("🔄 Loading embedding model...")
# embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
# print("✅ Embedding model loaded.")
#
# # === ENSURE NEW MODEL ===
# def ensure_qwen_model():
#     if not os.path.exists(MODEL_PATH):
#         print("📥 Downloading Qwen explanation model...")
#         with requests.get(MODEL_URL, stream=True) as r:
#             with open(MODEL_PATH, "wb") as f:
#                 for chunk in r.iter_content(1024 * 1024):
#                     f.write(chunk)
#         print("✔ Download complete.")
#     else:
#         print("✔ Qwen model exists.")
#
# def load_qwen():
#     print("⏳ Loading Qwen model...")
#     llm = GPT4All(model_name=MODEL_FILE, model_path=MODEL_FOLDER)
#     print("✔ Qwen model loaded.")
#     return llm
#
# ensure_qwen_model()
# qwen = load_qwen()
#
# # === UTILITIES ===
# def clean_path_input(s: str) -> str:
#     return s.strip().replace("\u202a", "").replace("\u202b", "").replace("\ufeff", "")
#
# def read_text_from_file(file_path: str) -> str:
#     ext = Path(file_path).suffix.lower()
#     text = ""
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             text = "\n".join([page.get_text("text") for page in doc])
#             doc.close()
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#         elif ext == ".pptx":
#             text = extract_text_from_pptx(file_path)
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#         elif ext in (".html", ".htm"):
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#         else:
#             raise ValueError(f"Unsupported file type: {ext}")
#     except Exception as e:
#         raise RuntimeError(f"Error reading {file_path}: {e}")
#     return text or ""
#
# def extract_text_from_pptx(pptx_path: str) -> str:
#     prs = Presentation(pptx_path)
#     all_text = []
#     for idx, slide in enumerate(prs.slides, start=1):
#         slide_text = []
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 slide_text.append(shape.text.strip())
#             try:
#                 if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
#                     blob = shape.image.blob
#                     img = Image.open(io.BytesIO(blob)).convert("RGB")
#                     ocr_text = pytesseract.image_to_string(img)
#                     if ocr_text.strip():
#                         slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
#             except:
#                 pass
#         if slide_text:
#             all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
#     return "\n\n".join(all_text)
#
# URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)
#
# def fetch_url_snippet(url: str) -> str:
#     headers = {"User-Agent": "Mozilla/5.0"}
#     try:
#         r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
#         r.raise_for_status()
#         soup = BeautifulSoup(r.text, "html.parser")
#         title = (soup.title.string.strip() if soup.title else "")
#         paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
#         snippet = (title + "\n\n" + " ".join(paragraphs))[:URL_SNIPPET_CHARS]
#         safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
#         with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
#             f.write(r.text)
#         return snippet or title or url
#     except:
#         return url
#
# def expand_links_in_text(text: str) -> str:
#     urls = list(dict.fromkeys(URL_RE.findall(text)))
#     if not urls:
#         return text
#     expanded_text = text
#     for url in urls:
#         snippet = fetch_url_snippet(url)
#         replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
#         expanded_text = expanded_text.replace(url, replacement)
#         sleep(0.2)
#     return expanded_text
#
# def clean_text(text: str) -> str:
#     text = re.sub(r'\r\n', '\n', text)
#     text = re.sub(r'\n{3,}', '\n\n', text)
#     text = re.sub(r'[ \t]+', ' ', text)
#     return text.strip()
#
# # === NEW EXPLANATION USING QWEN ===
# def chunk_words(text, size=CHUNK_SIZE):
#     words = text.split()
#     for i in range(0, len(words), size):
#         yield " ".join(words[i:i + size])
#
# def explain_chunk(chunk):
#     prompt = DEEP_PROMPT.format(chunk=chunk)
#     resp = qwen.generate(prompt, max_tokens=500)
#     return resp.strip()
#
# def explain_equation_general(equation):
#     prompt = EQUATION_PROMPT.format(equation=equation)
#     resp = qwen.generate(prompt, max_tokens=400)
#     return resp.strip()
#
# def generate_deep_explanation_with_equations_general(text):
#     final = []
#     parts = EQUATION_RE.split(text)
#     for part in parts:
#         if EQUATION_RE.match(part):
#             final.append(explain_equation_general(part))
#         else:
#             for chunk in chunk_words(part):
#                 final.append(explain_chunk(chunk))
#     return "\n\n".join(final)
#
# # === DATABASE / EMBEDDINGS ===
# def init_db():
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     c.execute("""
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB,
#             source_file TEXT
#         )
#     """)
#     conn.commit()
#     conn.close()
#
# def save_embeddings(source_file, chunks, vectors):
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     for ch, vec in zip(chunks, vectors):
#         c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
#                   (source_file, ch, vec.tobytes()))
#     conn.commit()
#     conn.close()
#
# # === AUDIO ===
# def save_english_audio(text, out_path):
#     try:
#         tts = pyttsx3.init()
#         tts.setProperty('rate', 150)
#         tts.setProperty('volume', 1.0)
#         tts.save_to_file(text, out_path)
#         tts.runAndWait()
#     except Exception as e:
#         print("❌ English audio failed:", e)
#
# def save_urdu_audio(text, out_path):
#     try:
#         gTTS(text=text.replace("\n", " "), lang='ur').save(out_path)
#     except Exception as e:
#         print("❌ Urdu audio failed:", e)
#
# # === SAFE URDU TRANSLATION ===
# def safe_translate_to_urdu(text, max_chunk=2500, retries=5, delay=2):
#     urdu_output = []
#     parts = [text[i:i+max_chunk] for i in range(0, len(text), max_chunk)]
#     for idx, part in enumerate(parts, start=1):
#         attempt = 1
#         while attempt <= retries:
#             try:
#                 print(f"🌐 Translating chunk {idx}/{len(parts)} (Attempt {attempt})")
#                 translated = GoogleTranslator(source='auto', target='ur').translate(part)
#                 urdu_output.append(translated)
#                 break
#             except Exception as e:
#                 print(f"⚠️ Translation error: {e}")
#                 attempt += 1
#                 time.sleep(delay)
#         if attempt > retries:
#             print("❌ Failed to translate this chunk after multiple retries.")
#             urdu_output.append(part)
#     return " ".join(urdu_output)
#
# # === MAIN PROCESS ===
# def process_document(file_path: str, expand_links: bool = True):
#     print(f"\n📄 Processing: {file_path}")
#
#     raw_text = read_text_from_file(file_path)
#
#     if expand_links:
#         raw_text = expand_links_in_text(raw_text)
#
#     cleaned = clean_text(raw_text)
#     base = Path(file_path).stem
#
#     # Save extracted text
#     raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
#     with open(raw_out, "w", encoding="utf-8") as f:
#         f.write(cleaned)
#
#     # Embeddings
#     chunks = [cleaned[i:i+1000] for i in range(0, len(cleaned), 1000)]
#     vectors = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(file_path, chunks, vectors)
#
#     # NEW DEEP EXPLANATION
#     print("🧠 Generating deep explanation using Qwen (including equations)...")
#     explanation = generate_deep_explanation_with_equations_general(cleaned)
#
#     explain_path = os.path.join(OUTPUT_DIR, f"{base}_deep_explanation.txt")
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation)
#
#     # English audio
#     save_english_audio(explanation, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))
#
#     # Urdu translation
#     ur_summary = safe_translate_to_urdu(explanation)
#
#     # Save Urdu text
#     ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
#     with open(ur_path, "w", encoding="utf-8") as f:
#         f.write(ur_summary)
#
#     # Urdu audio
#     save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))
#
#     print("\n✅ All outputs saved.")
#
# # === CLI ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
#     while True:
#         inp = input("👉 Enter file path: ").strip()
#         inp = clean_path_input(inp)
#         if inp.lower() in ("exit", "quit"):
#             break
#         if not os.path.exists(inp):
#             print("❌ File not found.")
#             continue
#         try:
#             process_document(inp, expand_links=True)
#         except Exception as e:
#             print("❌ Error:", e)
#

# working
# import os
# import io
# import re
# import sqlite3
# import requests
# import time
# from pathlib import Path
# from bs4 import BeautifulSoup
# from time import sleep
#
# # file-format libs
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation
# from PIL import Image
# import pandas as pd
#
# # NLP / tools
# import nltk
# nltk.download('punkt', quiet=True)
# from sentence_transformers import SentenceTransformer
# from deep_translator import GoogleTranslator
# import pyttsx3
# from gtts import gTTS
# import pytesseract
#
# # NEW MODEL (Qwen)
# from gpt4all import GPT4All
#
# # === CONFIG ===
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
# OUTPUT_DIR = "output"
# FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
# URL_FETCH_TIMEOUT = 10
# URL_SNIPPET_CHARS = 1200
#
# # NEW MODEL CONFIG
# MODEL_URL = "https://huggingface.co/Qwen/Qwen2.5-1.5B-Instruct-GGUF/resolve/main/qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_FILE = "qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_PATH = f"{MODEL_FOLDER}/{MODEL_FILE}"
# CHUNK_SIZE = 600
#
# # === IMPROVED DEEP EXPLANATION PROMPT ===
# DEEP_PROMPT = """
# You are an expert who provides deep and accurate explanations.
#
# You MUST follow these rules exactly:
#
# 1. Only explain the meaning of the given text.
# 2. Do NOT invent, add, assume, or hallucinate any details.
# 3. Do NOT repeat sentences from the text.
# 4. Do NOT summarize.
# 5. Do NOT fix the structure.
# 6. Only explain what exists inside the provided text.
# 7. If the text is unclear, explain ONLY what can be understood, do not guess.
# 8. Explanation must be a single coherent paragraph.
#
# TEXT STARTS:
# {chunk}
# TEXT ENDS.
# """
#
# # === IMPROVED EQUATION PROMPT ===
# EQUATION_PROMPT = """
# You are an expert at explaining equations from any field (math, physics, chemistry,
# economics, engineering, statistics).
#
# Strict Rules:
# 1. Identify WHAT TYPE of equation it is (if clear from symbols only).
# 2. Identify the FIELD of the equation (ONLY if inferable from structure).
# 3. Explain the PURPOSE of the equation only if it is directly clear from the symbols.
# 4. Explain the meaning of every variable that appears in the equation.
# 5. Explain the relationship between the variables.
# 6. Explain what the equation calculates or represents.
# 7. If it matches a known famous formula, mention the name (ONLY if structure clearly matches).
# 8. Do NOT hallucinate or assume any extra scenario.
# 9. Give a technical, detailed, structured explanation.
#
# Output format:
#
# ➤ Equation Detected:
# {equation}
#
# ➤ Type / Field:
# (...)
#
# ➤ Known Name (if clearly recognizable):
# (...)
#
# ➤ Variables Meaning:
# - x = ...
# - y = ...
# (Only variables that appear)
#
# ➤ What the Equation Describes:
# (...)
#
# ➤ How It Works:
# (...)
#
# ➤ Usage / Purpose:
# (Only if inferable)
#
# EQUATION START:
# {equation}
# EQUATION END.
# """
#
# # === IMPROVED EQUATION REGEX ===
# EQUATION_RE = re.compile(r'(\${1,2}[^$]+\${1,2})', flags=re.DOTALL)
#
# # Create folders
# os.makedirs(OUTPUT_DIR, exist_ok=True)
# os.makedirs(FETCHED_DIR, exist_ok=True)
# os.makedirs(MODEL_FOLDER, exist_ok=True)
#
# # === LOAD EMBEDDING MODEL ===
# print("🔄 Loading embedding model...")
# embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
# print("✅ Embedding model loaded.")
#
# # === ENSURE NEW MODEL ===
# def ensure_qwen_model():
#     if not os.path.exists(MODEL_PATH):
#         print("📥 Downloading Qwen explanation model...")
#         with requests.get(MODEL_URL, stream=True) as r:
#             with open(MODEL_PATH, "wb") as f:
#                 for chunk in r.iter_content(1024 * 1024):
#                     f.write(chunk)
#         print("✔ Download complete.")
#     else:
#         print("✔ Qwen model exists.")
#
# def load_qwen():
#     print("⏳ Loading Qwen model...")
#     llm = GPT4All(model_name=MODEL_FILE, model_path=MODEL_FOLDER)
#     print("✔ Qwen model loaded.")
#     return llm
#
# ensure_qwen_model()
# qwen = load_qwen()
#
# # === UTILITIES ===
# def clean_path_input(s: str) -> str:
#     return s.strip().replace("\u202a", "").replace("\u202b", "").replace("\ufeff", "")
#
# def read_text_from_file(file_path: str) -> str:
#     ext = Path(file_path).suffix.lower()
#     text = ""
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             text = "\n".join([page.get_text("text") for page in doc])
#             doc.close()
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#         elif ext == ".pptx":
#             text = extract_text_from_pptx(file_path)
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#         elif ext in (".html", ".htm"):
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#         else:
#             raise ValueError(f"Unsupported file type: {ext}")
#     except Exception as e:
#         raise RuntimeError(f"Error reading {file_path}: {e}")
#     return text or ""
#
# def extract_text_from_pptx(pptx_path: str) -> str:
#     prs = Presentation(pptx_path)
#     all_text = []
#     for idx, slide in enumerate(prs.slides, start=1):
#         slide_text = []
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 slide_text.append(shape.text.strip())
#             try:
#                 if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
#                     blob = shape.image.blob
#                     img = Image.open(io.BytesIO(blob)).convert("RGB")
#                     ocr_text = pytesseract.image_to_string(img)
#                     if ocr_text.strip():
#                         slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
#             except:
#                 pass
#         if slide_text:
#             all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
#     return "\n\n".join(all_text)
#
# URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)
#
# def fetch_url_snippet(url: str) -> str:
#     headers = {"User-Agent": "Mozilla/5.0"}
#     try:
#         r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
#         r.raise_for_status()
#         soup = BeautifulSoup(r.text, "html.parser")
#         title = (soup.title.string.strip() if soup.title else "")
#         paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
#         snippet = (title + "\n\n" + " ".join(paragraphs))[:URL_SNIPPET_CHARS]
#         safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
#         with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
#             f.write(r.text)
#         return snippet or title or url
#     except:
#         return url
#
# def expand_links_in_text(text: str) -> str:
#     urls = list(dict.fromkeys(URL_RE.findall(text)))
#     if not urls:
#         return text
#     expanded_text = text
#     for url in urls:
#         snippet = fetch_url_snippet(url)
#         replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
#         expanded_text = expanded_text.replace(url, replacement)
#         sleep(0.2)
#     return expanded_text
#
# def clean_text(text: str) -> str:
#     text = re.sub(r'\r\n', '\n', text)
#     text = re.sub(r'\n{3,}', '\n\n', text)
#     text = re.sub(r'[ \t]+', ' ', text)
#     return text.strip()
#
# # === NEW QWEN EXPLANATION SYSTEM ===
# def chunk_words(text, size=CHUNK_SIZE):
#     words = text.split()
#     for i in range(0, len(words), size):
#         yield " ".join(words[i:i + size])
#
# def explain_chunk(chunk):
#     prompt = DEEP_PROMPT.format(chunk=chunk)
#     resp = qwen.generate(prompt, max_tokens=500)
#     return resp.strip()
#
# def explain_equation_general(equation):
#     eq = equation.strip().replace("\n", " ")
#     prompt = EQUATION_PROMPT.format(equation=eq)
#     resp = qwen.generate(prompt, max_tokens=650)
#     return resp.strip()
#
# def generate_deep_explanation_with_equations_general(text):
#     final = []
#     parts = EQUATION_RE.split(text)
#
#     for part in parts:
#         if EQUATION_RE.match(part):
#             final.append("\n\n=== EQUATION EXPLANATION ===\n")
#             final.append(explain_equation_general(part))
#         else:
#             for chunk in chunk_words(part):
#                 final.append(explain_chunk(chunk))
#
#     return "\n\n".join(final)
#
# # === DATABASE / EMBEDDINGS ===
# def init_db():
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     c.execute("""
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB,
#             source_file TEXT
#         )
#     """)
#     conn.commit()
#     conn.close()
#
# def save_embeddings(source_file, chunks, vectors):
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     for ch, vec in zip(chunks, vectors):
#         c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
#                   (source_file, ch, vec.tobytes()))
#     conn.commit()
#     conn.close()
#
# # === AUDIO ===
# def save_english_audio(text, out_path):
#     try:
#         tts = pyttsx3.init()
#         tts.setProperty('rate', 150)
#         tts.setProperty('volume', 1.0)
#         tts.save_to_file(text, out_path)
#         tts.runAndWait()
#     except Exception as e:
#         print("❌ English audio failed:", e)
#
# def save_urdu_audio(text, out_path):
#     try:
#         gTTS(text=text.replace("\n", " "), lang='ur').save(out_path)
#     except Exception as e:
#         print("❌ Urdu audio failed:", e)
#
# # === SAFE URDU TRANSLATION ===
# def safe_translate_to_urdu(text, max_chunk=2500, retries=5, delay=2):
#     urdu_output = []
#     parts = [text[i:i+max_chunk] for i in range(0, len(text), max_chunk)]
#     for idx, part in enumerate(parts, start=1):
#         attempt = 1
#         while attempt <= retries:
#             try:
#                 print(f"🌐 Translating chunk {idx}/{len(parts)} (Attempt {attempt})")
#                 translated = GoogleTranslator(source='auto', target='ur').translate(part)
#                 urdu_output.append(translated)
#                 break
#             except Exception as e:
#                 print(f"⚠️ Translation error: {e}")
#                 attempt += 1
#                 time.sleep(delay)
#         if attempt > retries:
#             print("❌ Failed to translate this chunk after multiple retries.")
#             urdu_output.append(part)
#     return " ".join(urdu_output)
#
# # === MAIN PROCESS ===
# def process_document(file_path: str, expand_links: bool = True):
#     print(f"\n📄 Processing: {file_path}")
#
#     raw_text = read_text_from_file(file_path)
#
#     if expand_links:
#         raw_text = expand_links_in_text(raw_text)
#
#     cleaned = clean_text(raw_text)
#     base = Path(file_path).stem
#
#     # Save extracted text
#     raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
#     with open(raw_out, "w", encoding="utf-8") as f:
#         f.write(cleaned)
#
#     # Embeddings
#     chunks = [cleaned[i:i+1000] for i in range(0, len(cleaned), 1000)]
#     vectors = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(file_path, chunks, vectors)
#
#     # NEW DEEP EXPLANATION WITH EQUATIONS
#     print("🧠 Generating deep explanation using Qwen (including equations)...")
#     explanation = generate_deep_explanation_with_equations_general(cleaned)
#
#     explain_path = os.path.join(OUTPUT_DIR, f"{base}_deep_explanation.txt")
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation)
#
#     # English audio
#     save_english_audio(explanation, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))
#
#     # Urdu translation
#     ur_summary = safe_translate_to_urdu(explanation)
#
#     # Save Urdu text
#     ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
#     with open(ur_path, "w", encoding="utf-8") as f:
#         f.write(ur_summary)
#
#     # Urdu audio
#     save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))
#
#     print("\n✅ All outputs saved.")
#
# # === CLI ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
#     while True:
#         inp = input("👉 Enter file path: ").strip()
#         inp = clean_path_input(inp)
#         if inp.lower() in ("exit", "quit"):
#             break
#         if not os.path.exists(inp):
#             print("❌ File not found.")
#             continue
#         try:
#             process_document(inp, expand_links=True)
#         except Exception as e:
#             print("❌ Error:", e)

# import os
# import io
# import re
# import sqlite3
# import requests
# import time
# from pathlib import Path
# from bs4 import BeautifulSoup
# from time import sleep
#
# # file-format libs
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation
# from PIL import Image
# import pandas as pd
#
# # NLP / tools
# import nltk
#
# nltk.download('punkt', quiet=True)
# from sentence_transformers import SentenceTransformer
# from deep_translator import GoogleTranslator
# import pyttsx3
# from gtts import gTTS
# import pytesseract
#
# # NEW MODEL (your model)
# from gpt4all import GPT4All
#
# # === CONFIG ===
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
# OUTPUT_DIR = "output"
# FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
# URL_FETCH_TIMEOUT = 10
# URL_SNIPPET_CHARS = 1200
#
# # NEW MODEL CONFIG
# MODEL_URL = "https://huggingface.co/Qwen/Qwen2.5-1.5B-Instruct-GGUF/resolve/main/qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_FILE = "qwen2.5-1.5b-instruct-q4_0.gguf"
# MODEL_PATH = f"{MODEL_FOLDER}/{MODEL_FILE}"
# CHUNK_SIZE = 600
#
# DEEP_PROMPT = """
# You are an expert who provides deep and accurate explanations.
#
# You MUST follow these rules exactly:
#
# 1. Only explain the meaning of the given text.
# 2. Do NOT invent, add, assume, or hallucinate any details.
# 3. Do NOT repeat sentences from the text.
# 4. Do NOT summarize.
# 5. Do NOT fix the structure.
# 6. Only explain what exists inside the provided text.
# 7. If the text is unclear, explain ONLY what can be understood, do not guess.
# 8. Explanation must be a single coherent paragraph.
#
# TEXT STARTS:
# {chunk}
# TEXT ENDS.
# """
#
# # IMPROVED EQUATION PROMPT
# EQUATION_PROMPT = """
# You are an expert in explaining mathematical and scientific equations. For the given equation, provide a comprehensive explanation covering:
#
# 1. **Equation Identification**: What is this equation called? Which field does it belong to?
# 2. **Origin/Creator**: Who discovered/invented this equation (if known)?
# 3. **Symbol/Variable Meaning**: Explain each variable/symbol in the equation
# 4. **Physical/Mathematical Meaning**: What does this equation represent or calculate?
# 5. **Usage/Applications**: Where is this equation commonly used?
# 6. **Context**: What concepts or principles does it demonstrate?
#
# Provide detailed, accurate information. If you cannot identify the equation exactly, explain its structure and possible meanings based on the symbols used.
#
# EQUATION START:
# {equation}
# EQUATION END.
# """
#
# # IMPROVED EQUATION DETECTION - More comprehensive patterns
# EQUATION_PATTERNS = [
#     # LaTeX inline and block
#     r'\$\$(.*?)\$\$',
#     r'\$(.*?)\$',
#     # Mathematical expressions with common operators
#     r'[a-zA-Zα-ωΑ-Ω]_?[a-zA-Z0-9]*\s*=\s*[^=]+?(?=\n|$|\.)',
#     r'[a-zA-Zα-ωΑ-Ω]_?[a-zA-Z0-9]*\s*:\s*[^:]+?(?=\n|$|\.)',
#     # Common equation patterns
#     r'[a-zA-Z]+\s*\([^)]+\)\s*=\s*[^=]+?(?=\n|$|\.)',
#     # Chemical equations
#     r'[A-Z][a-z]?\d*\+?[A-Z][a-z]?\d*\-?→?[A-Z][a-z]?\d*\+?[A-Z][a-z]?\d*',
#     # Physics equations with Greek letters
#     r'[α-ωΑ-Ω]_?[a-zA-Z0-9]*\s*=\s*[^=]+?(?=\n|$|\.)',
# ]
#
# EQUATION_RE = re.compile('|'.join(EQUATION_PATTERNS), re.DOTALL)
#
# # Create folders
# os.makedirs(OUTPUT_DIR, exist_ok=True)
# os.makedirs(FETCHED_DIR, exist_ok=True)
# os.makedirs(MODEL_FOLDER, exist_ok=True)
#
# # === LOAD EMBEDDING MODEL ===
# print("🔄 Loading embedding model...")
# embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
# print("✅ Embedding model loaded.")
#
#
# # === ENSURE NEW MODEL ===
# def ensure_qwen_model():
#     if not os.path.exists(MODEL_PATH):
#         print("📥 Downloading Qwen explanation model...")
#         with requests.get(MODEL_URL, stream=True) as r:
#             with open(MODEL_PATH, "wb") as f:
#                 for chunk in r.iter_content(1024 * 1024):
#                     f.write(chunk)
#         print("✔ Download complete.")
#     else:
#         print("✔ Qwen model exists.")
#
#
# def load_qwen():
#     print("⏳ Loading Qwen model...")
#     llm = GPT4All(model_name=MODEL_FILE, model_path=MODEL_FOLDER)
#     print("✔ Qwen model loaded.")
#     return llm
#
#
# ensure_qwen_model()
# qwen = load_qwen()
#
#
# # === UTILITIES ===
# def clean_path_input(s: str) -> str:
#     return s.strip().replace("\u202a", "").replace("\u202b", "").replace("\ufeff", "")
#
#
# def read_text_from_file(file_path: str) -> str:
#     ext = Path(file_path).suffix.lower()
#     text = ""
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             text = "\n".join([page.get_text("text") for page in doc])
#             doc.close()
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#         elif ext == ".pptx":
#             text = extract_text_from_pptx(file_path)
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#         elif ext in (".html", ".htm"):
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#         else:
#             raise ValueError(f"Unsupported file type: {ext}")
#     except Exception as e:
#         raise RuntimeError(f"Error reading {file_path}: {e}")
#     return text or ""
#
#
# def extract_text_from_pptx(pptx_path: str) -> str:
#     prs = Presentation(pptx_path)
#     all_text = []
#     for idx, slide in enumerate(prs.slides, start=1):
#         slide_text = []
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 slide_text.append(shape.text.strip())
#             try:
#                 if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
#                     blob = shape.image.blob
#                     img = Image.open(io.BytesIO(blob)).convert("RGB")
#                     ocr_text = pytesseract.image_to_string(img)
#                     if ocr_text.strip():
#                         slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
#             except:
#                 pass
#         if slide_text:
#             all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
#     return "\n\n".join(all_text)
#
#
# URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)
#
#
# def fetch_url_snippet(url: str) -> str:
#     headers = {"User-Agent": "Mozilla/5.0"}
#     try:
#         r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
#         r.raise_for_status()
#         soup = BeautifulSoup(r.text, "html.parser")
#         title = (soup.title.string.strip() if soup.title else "")
#         paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
#         snippet = (title + "\n\n" + " ".join(paragraphs))[:URL_SNIPPET_CHARS]
#         safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
#         with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
#             f.write(r.text)
#         return snippet or title or url
#     except:
#         return url
#
#
# def expand_links_in_text(text: str) -> str:
#     urls = list(dict.fromkeys(URL_RE.findall(text)))
#     if not urls:
#         return text
#     expanded_text = text
#     for url in urls:
#         snippet = fetch_url_snippet(url)
#         replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
#         expanded_text = expanded_text.replace(url, replacement)
#         sleep(0.2)
#     return expanded_text
#
#
# def clean_text(text: str) -> str:
#     text = re.sub(r'\r\n', '\n', text)
#     text = re.sub(r'\n{3,}', '\n\n', text)
#     text = re.sub(r'[ \t]+', ' ', text)
#     return text.strip()
#
#
# # === IMPROVED EQUATION HANDLING ===
# def extract_equations(text):
#     """Extract equations from text with better filtering"""
#     equations = []
#     for pattern in EQUATION_PATTERNS:
#         matches = re.finditer(pattern, text, re.DOTALL)
#         for match in matches:
#             equation = match.group().strip()
#             # Filter out false positives
#             if is_valid_equation(equation):
#                 equations.append(equation)
#     return list(set(equations))  # Remove duplicates
#
#
# def is_valid_equation(equation):
#     """Check if the extracted text is actually an equation"""
#     # Too short to be meaningful
#     if len(equation) < 3:
#         return False
#
#     # Common non-equation patterns
#     non_equation_patterns = [
#         r'^[A-Z][a-z]+$',  # Single words
#         r'^\d+$',  # Just numbers
#         r'^[A-Za-z]+\s+[A-Za-z]+$',  # Two words
#     ]
#
#     for pattern in non_equation_patterns:
#         if re.match(pattern, equation.strip()):
#             return False
#
#     # Should contain mathematical elements
#     math_elements = ['=', '+', '-', '*', '/', '^', '(', ')', '{', '}', '[', ']', '\\']
#     if any(element in equation for element in math_elements):
#         return True
#
#     # Or Greek letters (common in equations)
#     greek_letters = ['α', 'β', 'γ', 'δ', 'ε', 'ζ', 'η', 'θ', 'ι', 'κ', 'λ', 'μ',
#                      'ν', 'ξ', 'π', 'ρ', 'σ', 'τ', 'υ', 'φ', 'χ', 'ψ', 'ω',
#                      'Α', 'Β', 'Γ', 'Δ', 'Ε', 'Ζ', 'Η', 'Θ', 'Ι', 'Κ', 'Λ', 'Μ',
#                      'Ν', 'Ξ', 'Π', 'Ρ', 'Σ', 'Τ', 'Υ', 'Φ', 'Χ', 'Ψ', 'Ω']
#
#     if any(letter in equation for letter in greek_letters):
#         return True
#
#     return False
#
#
# def explain_equation_comprehensive(equation):
#     """Generate comprehensive explanation for equations"""
#     prompt = EQUATION_PROMPT.format(equation=equation)
#     try:
#         resp = qwen.generate(prompt, max_tokens=600)
#         return resp.strip()
#     except Exception as e:
#         return f"Equation detected but explanation failed: {str(e)}\nEquation: {equation}"
#
#
# # === NEW EXPLANATION USING QWEN ===
# def chunk_words(text, size=CHUNK_SIZE):
#     words = text.split()
#     for i in range(0, len(words), size):
#         yield " ".join(words[i:i + size])
#
#
# def explain_chunk(chunk):
#     prompt = DEEP_PROMPT.format(chunk=chunk)
#     resp = qwen.generate(prompt, max_tokens=500)
#     return resp.strip()
#
#
# def generate_deep_explanation_with_equations_improved(text):
#     """Improved explanation with better equation handling"""
#     final = []
#
#     # First, extract all equations
#     equations = extract_equations(text)
#
#     if equations:
#         final.append("🔬 **EQUATIONS FOUND IN DOCUMENT:**")
#         final.append("=" * 50)
#
#         for i, eq in enumerate(equations, 1):
#             final.append(f"\n**Equation {i}:** `{eq}`")
#             explanation = explain_equation_comprehensive(eq)
#             final.append(f"**Explanation:** {explanation}")
#             final.append("-" * 30)
#
#     # Remove equations from text for regular explanation
#     text_without_equations = text
#     for eq in equations:
#         text_without_equations = text_without_equations.replace(eq, "")
#
#     # Generate regular explanation for non-equation text
#     if text_without_equations.strip():
#         final.append("\n📝 **DOCUMENT EXPLANATION:**")
#         final.append("=" * 50)
#         for chunk in chunk_words(text_without_equations):
#             if chunk.strip():
#                 explanation = explain_chunk(chunk)
#                 final.append(explanation)
#                 final.append("")  # Add spacing between chunks
#
#     return "\n".join(final)
#
#
# # === DATABASE / EMBEDDINGS ===
# def init_db():
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     c.execute("""
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB,
#             source_file TEXT
#         )
#     """)
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(source_file, chunks, vectors):
#     conn = sqlite3.connect(DB_PATH)
#     c = conn.cursor()
#     for ch, vec in zip(chunks, vectors):
#         c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
#                   (source_file, ch, vec.tobytes()))
#     conn.commit()
#     conn.close()
#
#
# # === AUDIO ===
# def save_english_audio(text, out_path):
#     try:
#         tts = pyttsx3.init()
#         tts.setProperty('rate', 150)
#         tts.setProperty('volume', 1.0)
#         tts.save_to_file(text, out_path)
#         tts.runAndWait()
#     except Exception as e:
#         print("❌ English audio failed:", e)
#
#
# def save_urdu_audio(text, out_path):
#     try:
#         gTTS(text=text.replace("\n", " "), lang='ur').save(out_path)
#     except Exception as e:
#         print("❌ Urdu audio failed:", e)
#
#
# # === SAFE URDU TRANSLATION ===
# def safe_translate_to_urdu(text, max_chunk=2500, retries=5, delay=2):
#     urdu_output = []
#     parts = [text[i:i + max_chunk] for i in range(0, len(text), max_chunk)]
#     for idx, part in enumerate(parts, start=1):
#         attempt = 1
#         while attempt <= retries:
#             try:
#                 print(f"🌐 Translating chunk {idx}/{len(parts)} (Attempt {attempt})")
#                 translated = GoogleTranslator(source='auto', target='ur').translate(part)
#                 urdu_output.append(translated)
#                 break
#             except Exception as e:
#                 print(f"⚠️ Translation error: {e}")
#                 attempt += 1
#                 time.sleep(delay)
#         if attempt > retries:
#             print("❌ Failed to translate this chunk after multiple retries.")
#             urdu_output.append(part)
#     return " ".join(urdu_output)
#
#
# # === MAIN PROCESS ===
# def process_document(file_path: str, expand_links: bool = True):
#     print(f"\n📄 Processing: {file_path}")
#
#     raw_text = read_text_from_file(file_path)
#
#     if expand_links:
#         raw_text = expand_links_in_text(raw_text)
#
#     cleaned = clean_text(raw_text)
#     base = Path(file_path).stem
#
#     # Save extracted text
#     raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
#     with open(raw_out, "w", encoding="utf-8") as f:
#         f.write(cleaned)
#
#     # Embeddings
#     chunks = [cleaned[i:i + 1000] for i in range(0, len(cleaned), 1000)]
#     vectors = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(file_path, chunks, vectors)
#
#     # IMPROVED DEEP EXPLANATION WITH BETTER EQUATION HANDLING
#     print("🧠 Generating deep explanation using Qwen (with improved equation detection)...")
#     explanation = generate_deep_explanation_with_equations_improved(cleaned)
#
#     explain_path = os.path.join(OUTPUT_DIR, f"{base}_deep_explanation.txt")
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation)
#
#     # English audio
#     save_english_audio(explanation, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))
#
#     # Urdu translation
#     ur_summary = safe_translate_to_urdu(explanation)
#
#     # Save Urdu text
#     ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
#     with open(ur_path, "w", encoding="utf-8") as f:
#         f.write(ur_summary)
#
#     # Urdu audio
#     save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))
#
#     print("\n✅ All outputs saved.")
#
#
# # === CLI ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
#     while True:
#         inp = input("👉 Enter file path: ").strip()
#         inp = clean_path_input(inp)
#         if inp.lower() in ("exit", "quit"):
#             break
#         if not os.path.exists(inp):
#             print("❌ File not found.")
#             continue
#         try:
#             process_document(inp, expand_links=True)
#         except Exception as e:
#             print("❌ Error:", e)

import os
import io
import re
import sqlite3
import requests
import time
from pathlib import Path
from bs4 import BeautifulSoup
from time import sleep

# File-format libraries
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pandas as pd

# NLP / Tools
import nltk
nltk.download('punkt', quiet=True)
from sentence_transformers import SentenceTransformer
from deep_translator import GoogleTranslator
import pyttsx3
from gtts import gTTS
import pytesseract

# New model
from gpt4all import GPT4All
from PyPDF2 import PdfReader

# === CONFIG ===
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
MODEL_FOLDER = "models"
DB_PATH = "embeddings.db"
OUTPUT_DIR = "output"
FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
URL_FETCH_TIMEOUT = 10
URL_SNIPPET_CHARS = 1200

# NEW MODEL CONFIG
MODEL_URL = "https://huggingface.co/Qwen/Qwen2.5-1.5B-Instruct-GGUF/resolve/main/qwen2.5-1.5b-instruct-q4_0.gguf"
MODEL_FILE = "qwen2.5-1.5b-instruct-q4_0.gguf"
MODEL_PATH = f"{MODEL_FOLDER}/{MODEL_FILE}"
CHUNK_SIZE = 600

# Equation-specific prompt
EQUATION_PROMPT = """
You are an expert who explains equations deeply.

Rules:
1. Only explain the equation provided.
2. Do NOT invent or summarize other content.
3. Give explanation in 2–3 lines.
4. If creator is not known, write 'Not specifically attributed'.

Now explain the following equation:

Equation: {equation}
"""

# General deep prompt for other text
DEEP_PROMPT = """
You are an expert who provides deep and accurate explanations.

You MUST follow these rules exactly:

1. Only explain the meaning of the given text.
2. Do NOT invent, add, assume, or hallucinate any details.
3. Do NOT repeat sentences from the text.
4. Do NOT summarize.
5. Do NOT fix the structure.
6. Only explain what exists inside the provided text.
7. If the text is unclear, explain ONLY what can be understood, do not guess.
8. Explanation must be a single coherent paragraph.

TEXT STARTS:
{chunk}
TEXT ENDS.
"""

# Create folders
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(FETCHED_DIR, exist_ok=True)
os.makedirs(MODEL_FOLDER, exist_ok=True)

# === LOAD EMBEDDING MODEL ===
print("🔄 Loading embedding model...")
embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
print("✅ Embedding model loaded.")

# === ENSURE & LOAD QWEN MODEL ===
def ensure_qwen_model():
    if not os.path.exists(MODEL_PATH):
        print("📥 Downloading Qwen explanation model...")
        with requests.get(MODEL_URL, stream=True) as r:
            with open(MODEL_PATH, "wb") as f:
                for chunk in r.iter_content(1024 * 1024):
                    f.write(chunk)
        print("✔ Download complete.")
    else:
        print("✔ Qwen model exists.")

def load_qwen():
    print("⏳ Loading Qwen model...")
    llm = GPT4All(model_name=MODEL_FILE, model_path=MODEL_FOLDER)
    print("✔ Qwen model loaded.")
    return llm

ensure_qwen_model()
qwen = load_qwen()

# === UTILITIES ===
def clean_path_input(s: str) -> str:
    # Remove invisible unicode characters and spaces
    invisible_chars = [
        '\u200b', '\u200c', '\u200d', '\u2060', '\u2061', '\u2062', '\u2063',
        '\u2064', '\u2066', '\u2067', '\u2068', '\u2069', '\u202a', '\u202b',
        '\u202c', '\u202d', '\u202e', '\ufeff'
    ]
    for char in invisible_chars:
        s = s.replace(char, "")
    return s.strip()

def read_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    text = ""
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            doc = fitz.open(file_path)
            text = "\n".join([page.get_text("text") for page in doc])
            doc.close()
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_path)
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif ext in (".html", ".htm"):
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
                text = soup.get_text(separator=" ")
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        raise RuntimeError(f"Error reading {file_path}: {e}")
    return text or ""

def extract_text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    all_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            try:
                if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
            except:
                pass
        if slide_text:
            all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
    return "\n\n".join(all_text)

URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)

def fetch_url_snippet(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        title = (soup.title.string.strip() if soup.title else "")
        paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
        snippet = (title + "\n\n" + " ".join(paragraphs))[:URL_SNIPPET_CHARS]
        safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
        with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
            f.write(r.text)
        return snippet or title or url
    except:
        return url

def expand_links_in_text(text: str) -> str:
    urls = list(dict.fromkeys(URL_RE.findall(text)))
    if not urls:
        return text
    expanded_text = text
    for url in urls:
        snippet = fetch_url_snippet(url)
        replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
        expanded_text = expanded_text.replace(url, replacement)
        sleep(0.2)
    return expanded_text

def clean_text(text: str) -> str:
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()

# === EQUATION FUNCTIONS ===
def read_pdf_for_equations(file_path):
    reader = PdfReader(file_path)
    text_pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        text = re.sub(r"[\n\r]+", " ", text)
        text = re.sub(r"\s{2,}", " ", text)
        text_pages.append(text.strip())
    return text_pages

def extract_equations_from_text(text):
    pattern = r"([A-Za-z0-9_\\^=∫\+\-\*/\(\)\[\]\,\.\s]+=[^,;]+|[A-Za-z0-9_\\^=∫\+\-\*/\(\)\[\]]+)"
    matches = re.findall(pattern, text)
    equations = [eq.strip() for eq in matches if len(eq.strip()) > 3]
    return equations

def explain_equation_qwen(model, equation):
    prompt = EQUATION_PROMPT.format(equation=equation)
    return model.generate(prompt, max_tokens=200).strip()

def process_equations(file_path, model):
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        pages = read_pdf_for_equations(file_path)
        content = " ".join(pages)
    else:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
    content = re.sub(r'\s+', ' ', content).strip()
    equations = extract_equations_from_text(content)
    print(f"✔ Found {len(equations)} equations.\n")
    final_output = ""
    for i, eq in enumerate(equations, start=1):
        print(f"🔍 Explaining Equation {i}: {eq}")
        explanation = explain_equation_qwen(model, eq)
        final_output += f"Equation {i}: {eq}\n{explanation}\n\n"
    out_path = os.path.join(OUTPUT_DIR, f"{Path(file_path).stem}_equation_explanations.txt")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(final_output)
    print(f"\n✅ Equation explanations saved to '{out_path}'")

# === DEEP EXPLANATION FUNCTIONS ===
def chunk_words(text, size=CHUNK_SIZE):
    words = text.split()
    for i in range(0, len(words), size):
        yield " ".join(words[i:i + size])

def explain_chunk(chunk):
    prompt = DEEP_PROMPT.format(chunk=chunk)
    resp = qwen.generate(prompt, max_tokens=500)
    return resp.strip()

def generate_deep_explanation(text):
    final = []
    for chunk in chunk_words(text):
        final.append(explain_chunk(chunk))
    return "\n\n".join(final)

# === DATABASE / EMBEDDINGS ===
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS embeddings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chunk TEXT,
            vector BLOB,
            source_file TEXT
        )
    """)
    conn.commit()
    conn.close()

def save_embeddings(source_file, chunks, vectors):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    for ch, vec in zip(chunks, vectors):
        c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
                  (source_file, ch, vec.tobytes()))
    conn.commit()
    conn.close()

# === AUDIO ===
def save_english_audio(text, out_path):
    try:
        tts = pyttsx3.init()
        tts.setProperty('rate', 150)
        tts.setProperty('volume', 1.0)
        tts.save_to_file(text, out_path)
        tts.runAndWait()
    except Exception as e:
        print("❌ English audio failed:", e)

def save_urdu_audio(text, out_path):
    try:
        gTTS(text=text.replace("\n", " "), lang='ur').save(out_path)
    except Exception as e:
        print("❌ Urdu audio failed:", e)

# === SAFE URDU TRANSLATION ===
def safe_translate_to_urdu(text, max_chunk=2500, retries=5, delay=2):
    urdu_output = []
    parts = [text[i:i+max_chunk] for i in range(0, len(text), max_chunk)]
    for idx, part in enumerate(parts, start=1):
        attempt = 1
        while attempt <= retries:
            try:
                print(f"🌐 Translating chunk {idx}/{len(parts)} (Attempt {attempt})")
                translated = GoogleTranslator(source='auto', target='ur').translate(part)
                urdu_output.append(translated)
                break
            except Exception as e:
                print(f"⚠️ Translation error: {e}")
                attempt += 1
                time.sleep(delay)
        if attempt > retries:
            print("❌ Failed to translate this chunk after multiple retries.")
            urdu_output.append(part)
    return " ".join(urdu_output)

# === MAIN DOCUMENT PROCESS ===
def process_document(file_path: str, expand_links: bool = True):
    print(f"\n📄 Processing: {file_path}")

    raw_text = read_text_from_file(file_path)

    if expand_links:
        raw_text = expand_links_in_text(raw_text)

    cleaned = clean_text(raw_text)
    base = Path(file_path).stem

    # Save extracted text
    raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
    with open(raw_out, "w", encoding="utf-8") as f:
        f.write(cleaned)

    # Embeddings
    chunks = [cleaned[i:i+1000] for i in range(0, len(cleaned), 1000)]
    vectors = embedder.encode(chunks, convert_to_numpy=True)
    save_embeddings(file_path, chunks, vectors)

    # NEW DEEP EXPLANATION
    print("🧠 Generating deep explanation using Qwen...")
    explanation = generate_deep_explanation(cleaned)
    explain_path = os.path.join(OUTPUT_DIR, f"{base}_deep_explanation.txt")
    with open(explain_path, "w", encoding="utf-8") as f:
        f.write(explanation)

    # English audio
    save_english_audio(explanation, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))

    # Urdu translation
    ur_summary = safe_translate_to_urdu(explanation)
    ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
    with open(ur_path, "w", encoding="utf-8") as f:
        f.write(ur_summary)
    save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))

    # === EQUATION PROCESSING ===
    print("📐 Extracting and explaining equations...")
    process_equations(file_path, qwen)

    print("\n✅ All outputs saved.")

# === CLI ===
if __name__ == "__main__":
    init_db()
    print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
    while True:
        inp = input("👉 Enter file path: ").strip()
        inp = clean_path_input(inp)
        if inp.lower() in ("exit", "quit"):
            break
        if not os.path.exists(inp):
            print("❌ File not found.")
            continue
        try:
            process_document(inp, expand_links=True)
        except Exception as e:
            print("❌ Error:", e)
