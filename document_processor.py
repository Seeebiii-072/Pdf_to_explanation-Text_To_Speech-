# """
# document_processor.py
# ----------------------
# Full document embedding pipeline:
# - Reads .txt or .pdf files
# - Cleans and splits text into chunks
# - Generates embeddings only
# - Saves them in SQLite for retrieval
# """
#
# import os
# import re
# import sqlite3
# from sentence_transformers import SentenceTransformer
# import fitz  # PyMuPDF
# import nltk
#
# nltk.download('punkt')
# from nltk.tokenize import sent_tokenize
#
# # === 1️⃣ Model Setup ===
# MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
# MODEL_FOLDER = "models"
#
# print("🔄 Loading embedding model...")
# model = SentenceTransformer(MODEL_NAME, cache_folder=MODEL_FOLDER)
# print("✅ Model loaded successfully!")
#
# # === 2️⃣ Helper Functions ===
#
# def read_text_from_file(file_path):
#     """Reads text from .txt or .pdf"""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#     if ext == ".txt":
#         with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#             text = f.read()
#     elif ext == ".pdf":
#         doc = fitz.open(file_path)
#         for page in doc:
#             text += page.get_text("text")
#         doc.close()
#     else:
#         raise ValueError("❌ Unsupported file format. Use .txt or .pdf")
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 3️⃣ Database Setup ===
#
# DB_PATH = "embeddings.db"
#
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute("INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#                        (text_chunk, emb.tobytes()))
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 4️⃣ Main Processing ===
#
# def process_document(file_path):
#     """Reads, cleans, embeds, and saves a document."""
#     print(f"📄 Reading document: {file_path}")
#     raw_text = read_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = model.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#     print("✅ Document processed and stored successfully!")
#
#
# # === 5️⃣ Run Example ===
# if __name__ == "__main__":
#     init_db()
#
#     # Example: place any .pdf or .txt in the same folder
#     FILE_PATH = r"C:\Users\Haseeb Ishtiaq\Desktop\AGREEMENT.pdf"  # change this to your test file
#     if os.path.exists(FILE_PATH):
#         process_document(FILE_PATH)
#     else:
#         print("⚠️ No document found. Please add a .pdf or .txt file to process.")


"""
document_processor.py
----------------------
Full document embedding + summarization pipeline:
- Reads .txt or .pdf files
- Cleans and splits text into chunks
- Generates embeddings
- Summarizes full text
- Saves embeddings in SQLite and summary file in same format
"""

"""
document_processor.py
----------------------
Full document embedding + summarization pipeline:
- Reads .txt or .pdf files
- Cleans and splits text into chunks
- Generates embeddings
- Summarizes full text
- Saves embeddings in SQLite and summary file in same format
"""
# hardcodeed good but not optimized
# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# print("🔄 Loading models (this may take a moment)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
# # === 2️⃣ Helper Functions ===
# def read_text_from_file(file_path):
#     """Reads text from .txt or .pdf"""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     if ext == ".txt":
#         with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#             text = f.read()
#     elif ext == ".pdf":
#         doc = fitz.open(file_path)
#         for page in doc:
#             text += page.get_text("text")
#         doc.close()
#     else:
#         raise ValueError("❌ Unsupported file format. Use .txt or .pdf")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 3️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 4️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     # If text too long, summarize in parts
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file with matching format."""
#     ext = os.path.splitext(file_path)[1].lower()
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary{ext}"
#
#     if ext == ".txt":
#         with open(summary_path, "w", encoding="utf-8") as f:
#             f.write(summary_text)
#     elif ext == ".pdf":
#         pdf = fitz.open()
#         page = pdf.new_page()
#         page.insert_text((50, 50), summary_text, fontsize=12)
#         pdf.save(summary_path)
#         pdf.close()
#     else:
#         print("⚠️ Unsupported format for summary saving.")
#         return
#
#     print(f"📝 Summary saved at: {summary_path}")
#
#
# # === 5️⃣ Main Processing ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, and saves a document."""
#     print(f"📄 Reading document: {file_path}")
#     raw_text = read_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     # Generate summary
#     summary = generate_summary(cleaned_text)
#     save_summary(file_path, summary)
#
#     print("✅ Document processed, embedded, and summarized successfully!")
#
#
# # === 6️⃣ Run Example ===
# if __name__ == "__main__":
#     init_db()
#
#     FILE_PATH = r"C:\Users\Haseeb Ishtiaq\Desktop\AGREEMENT.pdf"  # Change this path
#     if os.path.exists(FILE_PATH):
#         process_document(FILE_PATH)
#     else:
#         print("⚠️ No document found. Please add a .pdf or .txt file to process.")


"""
document_processor.py
----------------------
in loop and also optimized and till summary generator
Final optimized version:
- Loads embedding & summarization models once
- Loops to process multiple files (PDF or TXT)
- Reads, cleans, embeds, summarizes, and saves
- Stores embeddings in SQLite and summaries beside the originals
"""

# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# # === 2️⃣ Load Models Once ===
# print("🔄 Loading models (this may take a few moments)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!\n")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
#
# # === 3️⃣ Helper Functions ===
# def read_text_from_file(file_path):
#     """Reads text from .txt or .pdf"""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     if ext == ".txt":
#         with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#             text = f.read()
#     elif ext == ".pdf":
#         doc = fitz.open(file_path)
#         for page in doc:
#             text += page.get_text("text")
#         doc.close()
#     else:
#         raise ValueError("❌ Unsupported file format. Use .txt or .pdf")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 4️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 5️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file."""
#     ext = os.path.splitext(file_path)[1].lower()
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary{ext}"
#
#     if ext == ".txt":
#         with open(summary_path, "w", encoding="utf-8") as f:
#             f.write(summary_text)
#     elif ext == ".pdf":
#         pdf = fitz.open()
#         page = pdf.new_page()
#         page.insert_text((50, 50), summary_text, fontsize=12)
#         pdf.save(summary_path)
#         pdf.close()
#     else:
#         print("⚠️ Unsupported format for summary saving.")
#         return
#
#     print(f"📝 Summary saved at: {summary_path}")
#
#
# # === 6️⃣ Processing Logic ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, and saves a document."""
#     print(f"\n📄 Processing file: {file_path}")
#     raw_text = read_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     summary = generate_summary(cleaned_text)
#     save_summary(file_path, summary)
#
#     print("✅ Done! Document embedded and summarized.\n")
#
#
# # === 7️⃣ Main Loop ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Models loaded and database initialized.")
#     print("💡 Enter the full path of your .pdf or .txt file.")
#     print("🌀 Type 'exit' anytime to stop.\n")
#
#     while True:
#         file_path = input("👉 Enter file path: ").strip('"').strip()
#         if file_path.lower() in ["exit", "quit"]:
#             print("\n👋 Exiting program. All done!")
#             break
#
#         if os.path.exists(file_path):
#             try:
#                 process_document(file_path)
#             except Exception as e:
#                 print(f"❌ Error processing {file_path}: {e}")
#         else:
#             print("⚠️ File not found. Please check the path and try again.")
# API key
# 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

#explanation and summary done succesfull
# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# # === 2️⃣ Load Models Once ===
# print("🔄 Loading models (this may take a few moments)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!\n")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
#
# # === 3️⃣ Helper Functions ===
# def read_text_from_file(file_path):
#     """Reads text from .txt or .pdf"""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     if ext == ".txt":
#         with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#             text = f.read()
#     elif ext == ".pdf":
#         doc = fitz.open(file_path)
#         for page in doc:
#             text += page.get_text("text")
#         doc.close()
#     else:
#         raise ValueError("❌ Unsupported file format. Use .txt or .pdf")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 4️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 5️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file."""
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary.txt"
#
#     with open(summary_path, "w", encoding="utf-8") as f:
#         f.write(summary_text)
#
#     print(f"📝 Summary saved at: {summary_path}")
#     return summary_path
#
#
# # === 6️⃣ Explanation (using BART model) ===
# def generate_explanation(summary_text):
#     """Generate an easy-to-understand explanation using the same BART model."""
#     print("💬 Generating explanation (using BART)...")
#
#     prompt = (
#         "Explain the following text in simple and clear language so that anyone can easily understand it:\n\n"
#         f"{summary_text}\n\n"
#         "Keep it short, clear, and beginner-friendly."
#     )
#
#     explanation = summarizer(
#         prompt,
#         max_length=180,
#         min_length=60,
#         do_sample=True,
#         temperature=0.8,
#         top_p=0.9
#     )[0]['summary_text']
#
#     print("✅ Explanation generated successfully!\n")
#     return explanation
#
#
# def save_explanation(file_path, explanation_text):
#     """Save explanation as a .txt file next to summary."""
#     base_name = os.path.splitext(file_path)[0]
#     explain_path = f"{base_name}_explanation.txt"
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation_text)
#     print(f"📝 Explanation saved at: {explain_path}")
#
#
# # === 7️⃣ Processing Logic ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, and generates explanation."""
#     print(f"\n📄 Processing file: {file_path}")
#     raw_text = read_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     # Step 1: Summarize
#     summary = generate_summary(cleaned_text)
#     summary_path = save_summary(file_path, summary)
#
#     # Step 2: Explain via same model
#     explanation = generate_explanation(summary)
#     save_explanation(file_path, explanation)
#
#     print("✅ Document processed, summarized, and explained successfully!\n")
#
#
# # === 8️⃣ Main Loop ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Models loaded and database initialized.")
#     print("💡 Enter the full path of your .pdf or .txt file.")
#     print("🌀 Type 'exit' anytime to stop.\n")
#
#     while True:
#         file_path = input("👉 Enter file path: ").strip('"').strip()
#         if file_path.lower() in ["exit", "quit"]:
#             print("\n👋 Exiting program. All done!")
#             break
#
#         if os.path.exists(file_path):
#             try:
#                 process_document(file_path)
#             except Exception as e:
#                 print(f"❌ Error processing {file_path}: {e}")
#         else:
#             print("⚠️ File not found. Please check the path and try again.")

# full wprking with 2 file extension like pdf and txt
# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# from gtts import gTTS
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# # === 2️⃣ Load Models Once ===
# print("🔄 Loading models (this may take a few moments)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!\n")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
#
# # === 3️⃣ Helper Functions ===
# def read_text_from_file(file_path):
#     """Reads text from .txt or .pdf"""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     if ext == ".txt":
#         with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#             text = f.read()
#     elif ext == ".pdf":
#         doc = fitz.open(file_path)
#         for page in doc:
#             text += page.get_text("text")
#         doc.close()
#     else:
#         raise ValueError("❌ Unsupported file format. Use .txt or .pdf")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 4️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 5️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file."""
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary.txt"
#
#     with open(summary_path, "w", encoding="utf-8") as f:
#         f.write(summary_text)
#
#     print(f"📝 Summary saved at: {summary_path}")
#     return summary_path
#
#
# # === 6️⃣ Convert Summary to Audio ===
# def save_summary_audio(summary_text, file_path):
#     """Converts the generated summary into an audio file using gTTS."""
#     try:
#         print("\n🎧 Converting summary to audio...")
#         base_name = os.path.splitext(file_path)[0]
#         audio_folder = "audio_output"
#         os.makedirs(audio_folder, exist_ok=True)
#         audio_path = os.path.join(audio_folder, os.path.basename(base_name) + "_summary_audio.mp3")
#
#         tts = gTTS(text=summary_text, lang='en', slow=False)
#         tts.save(audio_path)
#
#         print(f"✅ Audio summary saved at: {audio_path}")
#         return audio_path
#     except Exception as e:
#         print(f"❌ Error converting summary to audio: {e}")
#
#
# # === 7️⃣ Explanation (using BART model) ===
# def generate_explanation(summary_text):
#     """Generate an easy-to-understand explanation using the same BART model."""
#     print("💬 Generating explanation (using BART)...")
#
#     prompt = (
#         "Explain the following text in simple and clear language so that anyone can easily understand it:\n\n"
#         f"{summary_text}\n\n"
#         "Keep it short, clear, and beginner-friendly."
#     )
#
#     explanation = summarizer(
#         prompt,
#         max_length=180,
#         min_length=60,
#         do_sample=True,
#         temperature=0.8,
#         top_p=0.9
#     )[0]['summary_text']
#
#     print("✅ Explanation generated successfully!\n")
#     return explanation
#
#
# def save_explanation(file_path, explanation_text):
#     """Save explanation as a .txt file next to summary."""
#     base_name = os.path.splitext(file_path)[0]
#     explain_path = f"{base_name}_explanation.txt"
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation_text)
#     print(f"📝 Explanation saved at: {explain_path}")
#
#
# # === 8️⃣ Processing Logic ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, generates explanation, and audio."""
#     print(f"\n📄 Processing file: {file_path}")
#     raw_text = read_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     # Step 1: Summarize
#     summary = generate_summary(cleaned_text)
#     summary_path = save_summary(file_path, summary)
#
#     # ✅ Step 1.5: Convert summary to audio
#     save_summary_audio(summary, file_path)
#
#     # Step 2: Explain via same model
#     explanation = generate_explanation(summary)
#     save_explanation(file_path, explanation)
#
#     print("✅ Document processed, summarized, explained, and audio generated successfully!\n")
#
#
# # === 9️⃣ Main Loop ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Models loaded and database initialized.")
#     print("💡 Enter the full path of your .pdf or .txt file.")
#     print("🌀 Type 'exit' anytime to stop.\n")
#
#     while True:
#         file_path = input("👉 Enter file path: ").strip('"').strip()
#         if file_path.lower() in ["exit", "quit"]:
#             print("\n👋 Exiting program. All done!")
#             break
#
#         if os.path.exists(file_path):
#             try:
#                 process_document(file_path)
#             except Exception as e:
#                 print(f"❌ Error processing {file_path}: {e}")
#         else:
#             print("⚠️ File not found. Please check the path and try again.")



# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# from gtts import gTTS
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
# from docx import Document
# from pptx import Presentation
# import pandas as pd
# from bs4 import BeautifulSoup
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# # === 2️⃣ Load Models Once ===
# print("🔄 Loading models (this may take a few moments)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!\n")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
#
# # === 3️⃣ Text Extraction Functions ===
# def extract_text_from_file(file_path):
#     """Reads text from pdf, docx, pptx, txt, csv, or html."""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             for page in doc:
#                 text += page.get_text("text")
#             doc.close()
#
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#
#         elif ext == ".pptx":
#             prs = Presentation(file_path)
#             slides_text = []
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         slides_text.append(shape.text)
#             text = "\n".join(slides_text)
#
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#
#         elif ext == ".html":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#
#         else:
#             raise ValueError("❌ Unsupported file type.")
#
#     except Exception as e:
#         raise RuntimeError(f"⚠️ Error reading {file_path}: {e}")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 4️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 5️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file."""
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary.txt"
#
#     with open(summary_path, "w", encoding="utf-8") as f:
#         f.write(summary_text)
#
#     print(f"📝 Summary saved at: {summary_path}")
#     return summary_path
#
#
# # === 6️⃣ Convert Summary to Audio ===
# def save_summary_audio(summary_text, file_path):
#     """Converts the generated summary into an audio file using gTTS."""
#     try:
#         print("\n🎧 Converting summary to audio...")
#         base_name = os.path.splitext(file_path)[0]
#         audio_folder = "audio_output"
#         os.makedirs(audio_folder, exist_ok=True)
#         audio_path = os.path.join(audio_folder, os.path.basename(base_name) + "_summary_audio.mp3")
#
#         tts = gTTS(text=summary_text, lang='en', slow=False)
#         tts.save(audio_path)
#
#         print(f"✅ Audio summary saved at: {audio_path}")
#         return audio_path
#     except Exception as e:
#         print(f"❌ Error converting summary to audio: {e}")
#
#
# # === 7️⃣ Explanation (using BART model) ===
# def generate_explanation(summary_text):
#     """Generate an easy-to-understand explanation using the same BART model."""
#     print("💬 Generating explanation (using BART)...")
#
#     prompt = (
#         "Explain the following text in  brief ,simple and clear language so that anyone can easily understand it:\n\n"
#         f"{summary_text}\n\n")
#
#     explanation = summarizer(
#         prompt,
#         max_length=180,
#         min_length=60,
#         do_sample=True,
#         temperature=0.8,
#         top_p=0.9
#     )[0]['summary_text']
#
#     print("✅ Explanation generated successfully!\n")
#     return explanation
#
#
# def save_explanation(file_path, explanation_text):
#     """Save explanation as a .txt file next to summary."""
#     base_name = os.path.splitext(file_path)[0]
#     explain_path = f"{base_name}_explanation.txt"
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation_text)
#     print(f"📝 Explanation saved at: {explain_path}")
#
#
# # === 8️⃣ Processing Logic ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, generates explanation, and audio."""
#     print(f"\n📄 Processing file: {file_path}")
#     raw_text = extract_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     # Step 1: Summarize
#     summary = generate_summary(cleaned_text)
#     summary_path = save_summary(file_path, summary)
#
#     # ✅ Step 1.5: Convert summary to audio
#     save_summary_audio(summary, file_path)
#
#     # Step 2: Explain via same model
#     explanation = generate_explanation(summary)
#     save_explanation(file_path, explanation)
#
#     print("✅ Document processed, summarized, explained, and audio generated successfully!\n")
#
#
# # === 9️⃣ Main Loop ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Models loaded and database initialized.")
#     print("💡 You can now process ANY file: .pdf, .docx, .pptx, .txt, .csv, .html, etc.")
#     print("🌀 Type 'exit' anytime to stop.\n")
#
#     while True:
#         file_path = input("👉 Enter file path: ").strip('"').strip()
#         if file_path.lower() in ["exit", "quit"]:
#             print("\n👋 Exiting program. All done!")
#             break
#
#         if os.path.exists(file_path):
#             try:
#                 process_document(file_path)
#             except Exception as e:
#                 print(f"❌ Error processing {file_path}: {e}")
#         else:
#             print("⚠️ File not found. Please check the path and try again.")

# without urdu and just text handling
# import os
# import re
# import sqlite3
# import fitz  # PyMuPDF
# import nltk
# import pyttsx3
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
# from docx import Document
# from pptx import Presentation
# import pandas as pd
# from bs4 import BeautifulSoup
#
# # === 1️⃣ Setup ===
# nltk.download('punkt', quiet=True)
# from nltk.tokenize import sent_tokenize
#
# EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
# SUMMARY_MODEL = "facebook/bart-large-cnn"
# MODEL_FOLDER = "models"
# DB_PATH = "embeddings.db"
#
# # === 2️⃣ Load Models Once ===
# print("🔄 Loading models (this may take a few moments)...")
# try:
#     embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
#     summarizer = pipeline("summarization", model=SUMMARY_MODEL)
#     print("✅ Models loaded successfully!\n")
# except Exception as e:
#     print(f"❌ Error loading models: {e}")
#     exit(1)
#
#
# # === 3️⃣ Text Extraction Functions ===
# def extract_text_from_file(file_path):
#     """Reads text from pdf, docx, pptx, txt, csv, or html."""
#     ext = os.path.splitext(file_path)[1].lower()
#     text = ""
#
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
#                 text = f.read()
#
#         elif ext == ".pdf":
#             doc = fitz.open(file_path)
#             for page in doc:
#                 text += page.get_text("text")
#             doc.close()
#
#         elif ext == ".docx":
#             doc = Document(file_path)
#             text = "\n".join([p.text for p in doc.paragraphs])
#
#         elif ext == ".pptx":
#             prs = Presentation(file_path)
#             slides_text = []
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         slides_text.append(shape.text)
#             text = "\n".join(slides_text)
#
#         elif ext == ".csv":
#             df = pd.read_csv(file_path)
#             text = df.to_string()
#
#         elif ext == ".html":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 soup = BeautifulSoup(f, "html.parser")
#                 text = soup.get_text(separator=" ")
#
#         else:
#             raise ValueError("❌ Unsupported file type.")
#
#     except Exception as e:
#         raise RuntimeError(f"⚠️ Error reading {file_path}: {e}")
#
#     return text
#
#
# def clean_text(text):
#     """Cleans text by removing extra spaces, newlines, and non-ASCII characters."""
#     text = re.sub(r'\s+', ' ', text)
#     text = re.sub(r'[^\x00-\x7F]+', ' ', text)
#     return text.strip()
#
#
# def split_into_chunks(text, max_sentences=5):
#     """Splits text into chunks of n sentences each for better embeddings."""
#     sentences = sent_tokenize(text)
#     chunks = [
#         " ".join(sentences[i:i + max_sentences])
#         for i in range(0, len(sentences), max_sentences)
#     ]
#     return chunks
#
#
# # === 4️⃣ Database Setup ===
# def init_db():
#     """Creates the embeddings table if not exists."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     cursor.execute('''
#         CREATE TABLE IF NOT EXISTS embeddings (
#             id INTEGER PRIMARY KEY AUTOINCREMENT,
#             chunk TEXT,
#             vector BLOB
#         )
#     ''')
#     conn.commit()
#     conn.close()
#
#
# def save_embeddings(chunks, embeddings):
#     """Saves chunks and their embeddings to the database."""
#     conn = sqlite3.connect(DB_PATH)
#     cursor = conn.cursor()
#     for text_chunk, emb in zip(chunks, embeddings):
#         cursor.execute(
#             "INSERT INTO embeddings (chunk, vector) VALUES (?, ?)",
#             (text_chunk, emb.tobytes())
#         )
#     conn.commit()
#     conn.close()
#     print(f"💾 Saved {len(chunks)} embeddings to database!")
#
#
# # === 5️⃣ Summarization ===
# def generate_summary(text):
#     """Generates a summary using the transformer model."""
#     print("🧠 Generating summary...")
#
#     max_len = 1024
#     if len(text) > max_len:
#         parts = [text[i:i + max_len] for i in range(0, len(text), max_len)]
#         summaries = summarizer(parts, max_length=150, min_length=40, do_sample=False)
#         combined = " ".join([s["summary_text"] for s in summaries])
#         return combined
#     else:
#         summary = summarizer(text, max_length=150, min_length=40, do_sample=False)
#         return summary[0]["summary_text"]
#
#
# def save_summary(file_path, summary_text):
#     """Saves the summary next to the original file."""
#     base_name = os.path.splitext(file_path)[0]
#     summary_path = f"{base_name}_summary.txt"
#
#     with open(summary_path, "w", encoding="utf-8") as f:
#         f.write(summary_text)
#
#     print(f"📝 Summary saved at: {summary_path}")
#     return summary_path
#
#
# # === 6️⃣ Convert Summary to Audio (Offline) ===
# def save_summary_audio(summary_text, file_path):
#     """Converts the generated summary into an offline audio file using pyttsx3."""
#     try:
#         print("\n🎧 Converting summary to audio (offline)...")
#         base_name = os.path.splitext(file_path)[0]
#         audio_folder = "audio_output"
#         os.makedirs(audio_folder, exist_ok=True)
#         audio_path = os.path.join(audio_folder, os.path.basename(base_name) + "_summary_audio.wav")
#
#         engine = pyttsx3.init()
#         engine.setProperty('rate', 160)   # speaking speed
#         engine.setProperty('volume', 1.0) # volume level
#
#         # Optional: switch voice (male/female)
#         voices = engine.getProperty('voices')
#         if voices:
#             engine.setProperty('voice', voices[0].id)  # change index if needed
#
#         engine.save_to_file(summary_text, audio_path)
#         engine.runAndWait()
#
#         print(f"✅ Offline audio summary saved at: {audio_path}")
#         return audio_path
#     except Exception as e:
#         print(f"❌ Error converting summary to audio: {e}")
#
#
# # === 7️⃣ Explanation (using BART model) ===
# def generate_explanation(summary_text):
#     """Generate an easy-to-understand explanation using the same BART model."""
#     print("💬 Generating explanation (using BART)...")
#
#     prompt = (
#         "Explain the following text in brief, simple and clear language so that anyone can easily understand it:\n\n"
#         f"{summary_text}\n\n")
#
#     explanation = summarizer(
#         prompt,
#         max_length=180,
#         min_length=60,
#         do_sample=True,
#         temperature=0.8,
#         top_p=0.9
#     )[0]['summary_text']
#
#     print("✅ Explanation generated successfully!\n")
#     return explanation
#
#
# def save_explanation(file_path, explanation_text):
#     """Save explanation as a .txt file next to summary."""
#     base_name = os.path.splitext(file_path)[0]
#     explain_path = f"{base_name}_explanation.txt"
#     with open(explain_path, "w", encoding="utf-8") as f:
#         f.write(explanation_text)
#     print(f"📝 Explanation saved at: {explain_path}")
#
#
# # === 8️⃣ Processing Logic ===
# def process_document(file_path):
#     """Reads, cleans, embeds, summarizes, generates explanation, and audio."""
#     print(f"\n📄 Processing file: {file_path}")
#     raw_text = extract_text_from_file(file_path)
#     cleaned_text = clean_text(raw_text)
#     chunks = split_into_chunks(cleaned_text)
#
#     print(f"🧩 Split into {len(chunks)} chunks. Generating embeddings...")
#     embeddings = embedder.encode(chunks, convert_to_numpy=True)
#     save_embeddings(chunks, embeddings)
#
#     # Step 1: Summarize
#     summary = generate_summary(cleaned_text)
#     save_summary(file_path, summary)
#
#     # ✅ Step 1.5: Convert summary to audio (offline)
#     save_summary_audio(summary, file_path)
#
#     # Step 2: Explain via same model
#     explanation = generate_explanation(summary)
#     save_explanation(file_path, explanation)
#
#     print("✅ Document processed, summarized, explained, and offline audio generated successfully!\n")
#
#
# # === 9️⃣ Main Loop ===
# if __name__ == "__main__":
#     init_db()
#     print("📂 Models loaded and database initialized.")
#     print("💡 You can now process ANY file: .pdf, .docx, .pptx, .txt, .csv, .html, etc.")
#     print("🌀 Type 'exit' anytime to stop.\n")
#
#     while True:
#         file_path = input("👉 Enter file path: ").strip('"').strip()
#         if file_path.lower() in ["exit", "quit"]:
#             print("\n👋 Exiting program. All done!")
#             break
#
#         if os.path.exists(file_path):
#             try:
#                 process_document(file_path)
#             except Exception as e:
#                 print(f"❌ Error processing {file_path}: {e}")
#         else:
#             print("⚠️ File not found. Please check the path and try again.")


#!/usr/bin/env python3 full working and good working

"""
unified_document_processor.py
- Multi-format text extraction (pdf, docx, pptx, txt, csv, html)
- PPTX image OCR (Tesseract)
- Active link expansion (fetch webpage content for URLs found)
- Summarization (facebook/bart-large-cnn)
- Explanation generation (same summarizer, short/simple)
- Embeddings (sentence-transformers) saved to SQLite
- English audio (pyttsx3, offline) + Urdu audio (gTTS, online)
- English summary + Urdu translation saved
"""

import os
import io
import re
import sqlite3
import requests
import numpy as np
from pathlib import Path
from bs4 import BeautifulSoup
from time import sleep

# file-format libs
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pandas as pd

# NLP / models / TTS
import nltk
nltk.download('punkt', quiet=True)
from nltk.tokenize import sent_tokenize

from sentence_transformers import SentenceTransformer
from transformers import pipeline
from deep_translator import GoogleTranslator
import pyttsx3
from gtts import gTTS
import pytesseract

# === CONFIG ===
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
SUMMARY_MODEL = "facebook/bart-large-cnn"
MODEL_FOLDER = "models"
DB_PATH = "embeddings.db"
OUTPUT_DIR = "output"
FETCHED_DIR = os.path.join(OUTPUT_DIR, "fetched_urls")
URL_FETCH_TIMEOUT = 10
URL_SNIPPET_CHARS = 1200

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(FETCHED_DIR, exist_ok=True)

# === LOAD MODELS ===
print("🔄 Loading models (this may take a moment)...")
try:
    embedder = SentenceTransformer(EMBED_MODEL, cache_folder=MODEL_FOLDER)
    summarizer = pipeline("summarization", model=SUMMARY_MODEL)
    print("✅ Models loaded.")
except Exception as e:
    print("❌ Error loading models:", e)
    raise

# === UTILITIES ===
def clean_path_input(s: str) -> str:
    return s.strip().replace("\u202a", "").replace("\u202b", "").replace("\ufeff", "")

def read_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    text = ""
    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            doc = fitz.open(file_path)
            text = "\n".join([page.get_text("text") for page in doc])
            doc.close()
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_path)
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif ext in (".html", ".htm"):
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
                text = soup.get_text(separator=" ")
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        raise RuntimeError(f"Error reading {file_path}: {e}")
    return text or ""

def extract_text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    all_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            try:
                if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        slide_text.append(f"[OCR image text]\n{ocr_text.strip()}")
            except Exception as e:
                print(f"⚠️ PPTX OCR error on slide {idx}: {e}")
        if slide_text:
            all_text.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
    return "\n\n".join(all_text)

URL_RE = re.compile(r'(https?://[^\s)>\]\}]+)', flags=re.IGNORECASE)

def fetch_url_snippet(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0 (compatible; DocumentProcessor/1.0)"}
    try:
        r = requests.get(url, headers=headers, timeout=URL_FETCH_TIMEOUT)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        title = (soup.title.string.strip() if soup.title else "") or ""
        paragraphs = [p.get_text(separator=" ", strip=True) for p in soup.find_all("p") if p.get_text(strip=True)]
        content = " ".join(paragraphs)
        snippet = (title + "\n\n" + content)[:URL_SNIPPET_CHARS].strip()
        safe_name = re.sub(r'[^0-9a-zA-Z_-]', '_', url)[:150]
        with open(os.path.join(FETCHED_DIR, f"{safe_name}.html"), "w", encoding="utf-8") as f:
            f.write(r.text)
        return snippet or title or url
    except Exception as e:
        print(f"⚠️ Failed to fetch {url}: {e}")
        return url

def expand_links_in_text(text: str) -> str:
    urls = list(dict.fromkeys(URL_RE.findall(text)))
    if not urls:
        return text
    print(f"🔗 Found {len(urls)} URL(s). Expanding...")
    expanded_text = text
    for url in urls:
        snippet = fetch_url_snippet(url)
        replacement = f"\n[Content expanded from {url}]\n{snippet}\n"
        expanded_text = expanded_text.replace(url, replacement)
        sleep(0.2)
    return expanded_text

def clean_text(text: str) -> str:
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()

# === ROBUST SUMMARIZATION ===
def summarize_text(text: str, chunk_size: int = 1000) -> str:
    if not text.strip():
        return "No content found."
    print("🧠 Running chunked summarization...")
    text = text.strip()
    summaries = []
    for i in range(0, len(text), chunk_size):
        chunk = text[i:i+chunk_size].strip()
        if not chunk:
            continue
        try:
            out = summarizer(chunk, max_length=150, min_length=40, do_sample=False)
            summaries.append(out[0]['summary_text'])
        except Exception as e:
            print(f"⚠️ Summarization failed for chunk {i//chunk_size + 1}: {e}")
            continue
    return " ".join(summaries) if summaries else "Summary could not be generated."

def generate_explanation(summary_text: str) -> str:
    if not summary_text.strip():
        return ""
    prompt = f"Explain the following text in short, simple and clear language:\n\n{summary_text}"
    try:
        out = summarizer(prompt, max_length=160, min_length=60, do_sample=False)
        return out[0]["summary_text"]
    except Exception as e:
        print(f"⚠️ Explanation generation failed: {e}")
        return summary_text

# === DATABASE / EMBEDDINGS ===
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS embeddings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            chunk TEXT,
            vector BLOB
        )
    """)
    # ensure 'source_file' exists
    c.execute("PRAGMA table_info(embeddings)")
    cols = [r[1] for r in c.fetchall()]
    if 'source_file' not in cols:
        print("⚙️ Upgrading database: adding 'source_file' column...")
        c.execute("ALTER TABLE embeddings ADD COLUMN source_file TEXT;")
        conn.commit()
        print("✅ Database upgraded.")
    conn.commit()
    conn.close()

def save_embeddings(source_file: str, chunks: list, vectors: np.ndarray):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    for ch, vec in zip(chunks, vectors):
        c.execute("INSERT INTO embeddings (source_file, chunk, vector) VALUES (?, ?, ?)",
                  (source_file, ch, vec.tobytes()))
    conn.commit()
    conn.close()
    print(f"💾 Saved {len(chunks)} embeddings for {source_file}")

def chunk_text_sentences(text: str, max_sentences=5):
    sents = sent_tokenize(text)
    chunks = [" ".join(sents[i:i+max_sentences]) for i in range(0, len(sents), max_sentences)]
    return [c for c in chunks if c.strip()]

# === AUDIO ===
#
def save_english_audio(text: str, out_path: str, voice_name="Zira"):
    try:
        print("🔊 Generating English audio ...")
        tts = pyttsx3.init()
        tts.setProperty('rate', 150)   # Speed
        tts.setProperty('volume', 1.0) # Volume

        # Select Zira or David
        for v in tts.getProperty('voices'):
            if voice_name.lower() in v.name.lower():
                tts.setProperty('voice', v.id)
                break

        tts.save_to_file(text, out_path)
        tts.runAndWait()
        print(f"✅ English audio saved at: {out_path}")
    except Exception as e:
        print(f"❌ English audio failed: {e}")

def save_urdu_audio(text: str, out_path: str):
    try:
        print("🔊 Generating Urdu audio ...")
        cleaned = text.replace("\n", " ").strip()
        gTTS(text=cleaned, lang='ur', slow=False).save(out_path)
        print(f"✅ Urdu audio saved at: {out_path}")
    except Exception as e:
        print(f"❌ Urdu audio failed: {e}")

# === MAIN PROCESS ===
def process_document(file_path: str, expand_links: bool = True):
    print(f"\n📄 Processing: {file_path}")
    raw_text = read_text_from_file(file_path)
    if not raw_text.strip():
        print("⚠️ No text extracted.")
        return

    if expand_links:
        raw_text = expand_links_in_text(raw_text)

    cleaned = clean_text(raw_text)

    base = Path(file_path).stem
    raw_out = os.path.join(OUTPUT_DIR, f"{base}_extracted.txt")
    with open(raw_out, "w", encoding="utf-8") as f:
        f.write(cleaned)
    print(f"📂 Raw extracted text saved: {raw_out}")

    chunks = chunk_text_sentences(cleaned, max_sentences=5)
    if chunks:
        print(f"🧩 Creating {len(chunks)} embedding chunks ...")
        vectors = embedder.encode(chunks, convert_to_numpy=True)
        save_embeddings(file_path, chunks, vectors)
    else:
        print("⚠️ No chunks for embeddings.")

    en_summary = summarize_text(cleaned)
    en_summary_path = os.path.join(OUTPUT_DIR, f"{base}_english_summary.txt")
    with open(en_summary_path, "w", encoding="utf-8") as f:
        f.write(en_summary)
    print(f"📝 English summary saved: {en_summary_path}")

    explanation = generate_explanation(en_summary)
    explain_path = os.path.join(OUTPUT_DIR, f"{base}_explanation.txt")
    with open(explain_path, "w", encoding="utf-8") as f:
        f.write(explanation)
    print(f"💡 Explanation saved: {explain_path}")

    save_english_audio(en_summary, os.path.join(OUTPUT_DIR, f"{base}_english_audio.mp3"))

    print("🌐 Translating summary to Urdu...")
    try:
        if len(en_summary) > 3500:
            parts = [en_summary[i:i+3500] for i in range(0, len(en_summary), 3500)]
            urd_parts = [GoogleTranslator(source='auto', target='ur').translate(p) for p in parts]
            ur_summary = " ".join(urd_parts)
        else:
            ur_summary = GoogleTranslator(source='auto', target='ur').translate(en_summary)
    except Exception as e:
        print(f"⚠️ Translation failed: {e}")
        ur_summary = en_summary

    ur_path = os.path.join(OUTPUT_DIR, f"{base}_urdu_summary.txt")
    with open(ur_path, "w", encoding="utf-8") as f:
        f.write(ur_summary)
    print(f"📝 Urdu summary saved: {ur_path}")

    save_urdu_audio(ur_summary, os.path.join(OUTPUT_DIR, f"{base}_urdu_audio.mp3"))

    print("\n✅ All outputs saved in 'output/'")

# === CLI ===
if __name__ == "__main__":
    init_db()
    print("📂 Ready. Enter file path (pdf, docx, pptx, txt, csv, html). Type 'exit' to quit.\n")
    while True:
        inp = input("👉 Enter file path: ").strip()
        inp = clean_path_input(inp)
        if inp.lower() in ("exit", "quit"):
            print("👋 Bye")
            break
        if not inp or not os.path.exists(inp):
            print(f"❌ File not found: {inp}")
            continue
        try:
            process_document(inp, expand_links=True)
        except Exception as e:
            print(f"❌ Processing failed: {e}")
